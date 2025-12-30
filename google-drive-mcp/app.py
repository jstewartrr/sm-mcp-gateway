"""
Google Drive MCP Server v3.3.0 - With full read/write support
"""
import os, json, io, uuid, base64
from flask import Flask, request, jsonify, make_response
from flask_cors import CORS
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

app = Flask(__name__)
CORS(app, resources={r"/mcp": {"origins": "*", "methods": ["POST", "OPTIONS"], "allow_headers": ["Content-Type", "Mcp-Session-Id"], "expose_headers": ["Mcp-Session-Id"]}})
sessions = {}
# Full access scope for read AND write
SCOPES = ['https://www.googleapis.com/auth/drive']

def get_drive_service():
    creds_json = os.environ.get('GOOGLE_SERVICE_ACCOUNT_JSON')
    if not creds_json: raise ValueError("GOOGLE_SERVICE_ACCOUNT_JSON not set")
    return build('drive', 'v3', credentials=service_account.Credentials.from_service_account_info(json.loads(creds_json), scopes=SCOPES))

TOOLS = [
    {"name": "list_shared_drives", "description": "List all Shared Drives accessible to the service account", "inputSchema": {"type": "object", "properties": {}, "required": []}},
    {"name": "list_folder_contents", "description": "List files in a folder (works with Shared Drives)", "inputSchema": {"type": "object", "properties": {"folder_id": {"type": "string", "description": "Folder ID or Shared Drive ID"}, "page_size": {"type": "integer", "default": 50}}, "required": ["folder_id"]}},
    {"name": "search_files", "description": "Search files by name (includes Shared Drives)", "inputSchema": {"type": "object", "properties": {"query": {"type": "string"}, "folder_id": {"type": "string"}, "file_type": {"type": "string"}}, "required": ["query"]}},
    {"name": "get_file_metadata", "description": "Get file metadata", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}}, "required": ["file_id"]}},
    {"name": "read_text_file", "description": "Read text files", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}}, "required": ["file_id"]}},
    {"name": "read_excel_file", "description": "Read Excel/Sheets", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}, "sheet_name": {"type": "string"}}, "required": ["file_id"]}},
    {"name": "read_pdf_file", "description": "Extract PDF text", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}, "page_numbers": {"type": "array", "items": {"type": "integer"}}}, "required": ["file_id"]}},
    {"name": "read_word_file", "description": "Extract Word text", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}}, "required": ["file_id"]}},
    {"name": "read_powerpoint_file", "description": "Extract PowerPoint text", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string"}}, "required": ["file_id"]}},
    {"name": "create_folder", "description": "Create a new folder in Google Drive or Shared Drive", "inputSchema": {"type": "object", "properties": {"folder_name": {"type": "string", "description": "Name of the folder to create"}, "parent_id": {"type": "string", "description": "Parent folder ID or Shared Drive ID (optional)"}}, "required": ["folder_name"]}},
    {"name": "upload_file", "description": "Upload a file to Google Drive (base64 encoded content)", "inputSchema": {"type": "object", "properties": {"file_name": {"type": "string", "description": "Name for the file"}, "content": {"type": "string", "description": "Base64 encoded file content"}, "folder_id": {"type": "string", "description": "Destination folder ID (optional)"}, "mime_type": {"type": "string", "description": "MIME type of the file (e.g., text/plain, application/pdf)"}}, "required": ["file_name", "content"]}},
    {"name": "move_file", "description": "Move a file to a different folder", "inputSchema": {"type": "object", "properties": {"file_id": {"type": "string", "description": "ID of the file to move"}, "new_parent_id": {"type": "string", "description": "ID of the destination folder"}}, "required": ["file_id", "new_parent_id"]}}
]

def download_file(svc, fid):
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, svc.files().get_media(fileId=fid, supportsAllDrives=True))
    done = False
    while not done: _, done = dl.next_chunk()
    buf.seek(0)
    return buf

def export_file(svc, fid, mime):
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, svc.files().export_media(fileId=fid, mimeType=mime))
    done = False
    while not done: _, done = dl.next_chunk()
    buf.seek(0)
    return buf

def list_shared_drives():
    svc = get_drive_service()
    drives = svc.drives().list(pageSize=100).execute().get('drives', [])
    return {"shared_drives": [{"id": d['id'], "name": d['name']} for d in drives]}

def list_folder_contents(folder_id, page_size=50):
    svc = get_drive_service()
    files = svc.files().list(
        q=f"'{folder_id}' in parents and trashed=false",
        pageSize=min(page_size, 100),
        fields="files(id,name,mimeType,size,modifiedTime)",
        supportsAllDrives=True,
        includeItemsFromAllDrives=True
    ).execute().get('files', [])
    return {"files": [{"id": f['id'], "name": f['name'], "type": f['mimeType'], "size": f.get('size', 'N/A'), "modified": f.get('modifiedTime')} for f in files]}

def search_files(query, folder_id=None, file_type=None):
    svc = get_drive_service()
    q = f"name contains '{query}' and trashed=false"
    if folder_id: q += f" and '{folder_id}' in parents"
    mime_map = {"spreadsheet": "application/vnd.google-apps.spreadsheet", "document": "application/vnd.google-apps.document", "pdf": "application/pdf", "presentation": "application/vnd.google-apps.presentation", "folder": "application/vnd.google-apps.folder"}
    if file_type in mime_map: q += f" and mimeType='{mime_map[file_type]}'"
    return {"files": svc.files().list(q=q, pageSize=50, fields="files(id,name,mimeType,size)", supportsAllDrives=True, includeItemsFromAllDrives=True).execute().get('files', [])}

def read_excel_file(file_id, sheet_name=None):
    svc = get_drive_service()
    meta = svc.files().get(fileId=file_id, fields="mimeType,name", supportsAllDrives=True).execute()
    buf = export_file(svc, file_id, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet') if 'spreadsheet' in meta.get('mimeType', '') else download_file(svc, file_id)
    if sheet_name:
        df = pd.read_excel(buf, sheet_name=sheet_name)
        return {"file_name": meta['name'], "sheet": sheet_name, "columns": list(df.columns), "row_count": len(df), "data": df.head(100).to_dict(orient='records')}
    xlsx = pd.ExcelFile(buf)
    return {"file_name": meta['name'], "sheets": xlsx.sheet_names, "data": {s: {"columns": list((df := pd.read_excel(xlsx, sheet_name=s)).columns), "rows": len(df), "data": df.head(50).to_dict(orient='records')} for s in xlsx.sheet_names[:5]}}

def read_pdf_file(file_id, page_numbers=None):
    svc = get_drive_service()
    meta = svc.files().get(fileId=file_id, fields="name", supportsAllDrives=True).execute()
    reader = PdfReader(download_file(svc, file_id))
    pages = page_numbers or range(min(len(reader.pages), 20))
    return {"file_name": meta['name'], "total_pages": len(reader.pages), "content": [{"page": p+1, "text": reader.pages[p].extract_text()} for p in pages if 0 <= p < len(reader.pages)]}

def read_powerpoint_file(file_id):
    svc = get_drive_service()
    meta = svc.files().get(fileId=file_id, fields="mimeType,name", supportsAllDrives=True).execute()
    buf = export_file(svc, file_id, 'application/vnd.openxmlformats-officedocument.presentationml.presentation') if 'presentation' in meta.get('mimeType', '') else download_file(svc, file_id)
    prs = Presentation(buf)
    return {"file_name": meta['name'], "slides": [{"num": i+1, "text": "\n".join([s.text for s in sl.shapes if hasattr(s, "text")])} for i, sl in enumerate(prs.slides)]}

def read_word_file(file_id):
    svc = get_drive_service()
    meta = svc.files().get(fileId=file_id, fields="mimeType,name", supportsAllDrives=True).execute()
    buf = export_file(svc, file_id, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document') if 'document' in meta.get('mimeType', '') else download_file(svc, file_id)
    return {"file_name": meta['name'], "paragraphs": [p.text for p in Document(buf).paragraphs if p.text.strip()][:100]}

def read_text_file(file_id):
    svc = get_drive_service()
    meta = svc.files().get(fileId=file_id, fields="name", supportsAllDrives=True).execute()
    return {"file_name": meta['name'], "content": download_file(svc, file_id).read().decode('utf-8', errors='replace')[:50000]}

def get_file_metadata(file_id):
    svc = get_drive_service()
    f = svc.files().get(fileId=file_id, fields="id,name,mimeType,size,createdTime,modifiedTime,owners,webViewLink,driveId", supportsAllDrives=True).execute()
    return {"id": f['id'], "name": f['name'], "type": f['mimeType'], "size": f.get('size'), "created": f.get('createdTime'), "modified": f.get('modifiedTime'), "owner": f.get('owners', [{}])[0].get('emailAddress') if f.get('owners') else None, "link": f.get('webViewLink'), "shared_drive_id": f.get('driveId')}

def create_folder(folder_name, parent_id=None):
    svc = get_drive_service()
    metadata = {'name': folder_name, 'mimeType': 'application/vnd.google-apps.folder'}
    if parent_id: metadata['parents'] = [parent_id]
    folder = svc.files().create(body=metadata, fields='id,name,webViewLink', supportsAllDrives=True).execute()
    return {"id": folder['id'], "name": folder['name'], "link": folder.get('webViewLink')}

def upload_file(file_name, content, folder_id=None, mime_type=None):
    svc = get_drive_service()
    try:
        file_content = base64.b64decode(content)
    except Exception as e:
        raise ValueError(f"Invalid base64 content: {str(e)}")
    if not mime_type:
        ext_map = {'.txt': 'text/plain', '.pdf': 'application/pdf', '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation', '.json': 'application/json', '.csv': 'text/csv', '.html': 'text/html', '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg'}
        mime_type = next((v for k, v in ext_map.items() if file_name.lower().endswith(k)), 'application/octet-stream')
    metadata = {'name': file_name}
    if folder_id: metadata['parents'] = [folder_id]
    media = MediaIoBaseUpload(io.BytesIO(file_content), mimetype=mime_type, resumable=True)
    file = svc.files().create(body=metadata, media_body=media, fields='id,name,webViewLink,size', supportsAllDrives=True).execute()
    return {"id": file['id'], "name": file['name'], "link": file.get('webViewLink'), "size": file.get('size')}

def move_file(file_id, new_parent_id):
    svc = get_drive_service()
    file = svc.files().get(fileId=file_id, fields='parents', supportsAllDrives=True).execute()
    previous_parents = ",".join(file.get('parents', []))
    file = svc.files().update(fileId=file_id, addParents=new_parent_id, removeParents=previous_parents, fields='id,name,parents,webViewLink', supportsAllDrives=True).execute()
    return {"id": file['id'], "name": file['name'], "new_parent": new_parent_id, "link": file.get('webViewLink')}

def execute_tool(name, args):
    dispatch = {
        "list_shared_drives": lambda: list_shared_drives(),
        "list_folder_contents": lambda: list_folder_contents(args.get('folder_id'), args.get('page_size', 50)),
        "search_files": lambda: search_files(args.get('query'), args.get('folder_id'), args.get('file_type')),
        "read_excel_file": lambda: read_excel_file(args.get('file_id'), args.get('sheet_name')),
        "read_pdf_file": lambda: read_pdf_file(args.get('file_id'), args.get('page_numbers')),
        "read_powerpoint_file": lambda: read_powerpoint_file(args.get('file_id')),
        "read_word_file": lambda: read_word_file(args.get('file_id')),
        "read_text_file": lambda: read_text_file(args.get('file_id')),
        "get_file_metadata": lambda: get_file_metadata(args.get('file_id')),
        "create_folder": lambda: create_folder(args.get('folder_name'), args.get('parent_id')),
        "upload_file": lambda: upload_file(args.get('file_name'), args.get('content'), args.get('folder_id'), args.get('mime_type')),
        "move_file": lambda: move_file(args.get('file_id'), args.get('new_parent_id'))
    }
    if name not in dispatch: raise ValueError(f"Unknown tool: {name}")
    return dispatch[name]()

@app.route('/mcp', methods=['POST', 'OPTIONS'])
def mcp():
    if request.method == 'OPTIONS':
        r = make_response('', 204)
        r.headers['Access-Control-Allow-Origin'] = '*'
        r.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        r.headers['Access-Control-Allow-Headers'] = 'Content-Type, Mcp-Session-Id'
        r.headers['Access-Control-Expose-Headers'] = 'Mcp-Session-Id'
        return r
    try: data = request.get_json()
    except: return jsonify({"jsonrpc": "2.0", "error": {"code": -32700, "message": "Parse error"}, "id": None})
    if not data: return jsonify({"jsonrpc": "2.0", "error": {"code": -32700, "message": "Parse error"}, "id": None})
    method, params, rid = data.get('method'), data.get('params', {}), data.get('id')
    sid = request.headers.get('Mcp-Session-Id')
    if not sid and method == "initialize": sid = str(uuid.uuid4()); sessions[sid] = True
    try:
        if method == 'initialize': result = {"protocolVersion": "2024-11-05", "serverInfo": {"name": "google-drive-mcp", "version": "3.3.0"}, "capabilities": {"tools": {}}}
        elif method == 'notifications/initialized': return '', 204
        elif method == 'tools/list': result = {"tools": TOOLS}
        elif method == 'tools/call': result = {"content": [{"type": "text", "text": json.dumps(execute_tool(params.get('name'), params.get('arguments', {})), indent=2, default=str)}]}
        else: return jsonify({"jsonrpc": "2.0", "error": {"code": -32601, "message": f"Unknown: {method}"}, "id": rid})
        r = make_response(jsonify({"jsonrpc": "2.0", "result": result, "id": rid}))
        r.headers['Access-Control-Expose-Headers'] = 'Mcp-Session-Id'
        if sid: r.headers['Mcp-Session-Id'] = sid
        return r
    except Exception as e: return jsonify({"jsonrpc": "2.0", "error": {"code": -32000, "message": str(e)}, "id": rid})

@app.route('/health')
def health(): return jsonify({"status": "healthy"})

@app.route('/')
def root(): return jsonify({"service": "google-drive-mcp", "status": "healthy", "tools": [t['name'] for t in TOOLS], "version": "3.3.0"})

if __name__ == '__main__': app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 8080)))
